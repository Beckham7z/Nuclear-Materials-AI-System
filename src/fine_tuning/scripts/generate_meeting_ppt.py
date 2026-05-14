#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
核材料垂类专家模型项目 - 组会汇报PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
TITLE_COLOR = RGBColor(0, 51, 102)  # 深蓝色
ACCENT_COLOR = RGBColor(0, 102, 204)  # 亮蓝色
TEXT_COLOR = RGBColor(51, 51, 51)  # 深灰色
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景形状
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, notes=""):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题区域背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = TITLE_COLOR
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """添加两栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题区域背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = TITLE_COLOR
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.text = left_title
    left_title_p.font.size = Pt(24)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = ACCENT_COLOR
    
    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.text = right_title
    right_title_p.font.size = Pt(24)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = ACCENT_COLOR
    
    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    return slide

def add_progress_slide(prs):
    """添加进展时间线"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题区域背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = TITLE_COLOR
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "项目进展时间线"
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # 时间线内容
    timeline_items = [
        ("第一阶段", "数据收集与清洗", "完成核材料领域文档收集、QA对生成、数据清洗"),
        ("第二阶段", "模型训练DAPT", "完成Qwen3.5-2B/9B模型微调，LoRA+SWIFT框架"),
        ("第三阶段", "RAG系统构建", "搭建知识图谱+向量数据库混合检索"),
        ("第四阶段", "Agent框架集成", "实现Plan-Act双模式推理框架"),
    ]
    
    for i, (phase, title, desc) in enumerate(timeline_items):
        y_pos = Inches(1.6 + i * 1.4)
        
        # 阶段标记
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y_pos + Inches(0.1), Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_COLOR
        circle.line.fill.background()
        
        # 阶段标题
        phase_box = slide.shapes.add_textbox(Inches(1.5), y_pos, Inches(2), Inches(0.4))
        phase_frame = phase_box.text_frame
        phase_p = phase_frame.paragraphs[0]
        phase_p.text = phase
        phase_p.font.size = Pt(16)
        phase_p.font.bold = True
        phase_p.font.color.rgb = ACCENT_COLOR
        
        # 标题和描述
        content_box = slide.shapes.add_textbox(Inches(3.5), y_pos, Inches(9), Inches(0.8))
        content_frame = content_box.text_frame
        title_p = content_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = TEXT_COLOR
        
        desc_p = content_frame.add_paragraph()
        desc_p.text = desc
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = TEXT_COLOR
    
    return slide

# ==================== 制作PPT ====================

# 1. 标题页
add_title_slide(prs, 
    "核材料垂类专家模型项目汇报",
    "AI4S: 面向核电领域的智能问答与知识推理系统")

# 2. 背景
add_content_slide(prs, "研究背景与动机", [
    "核电行业发展迅速，但专业知识获取效率低",
    "核材料领域知识专业性强、更新快，检索困难",
    "传统问答系统缺乏深度推理能力",
    "AI for Science (AI4S) 成为科研新范式",
    "大语言模型为领域专家系统带来新机遇"
])

# 3. 研究目标
add_content_slide(prs, "研究目标", [
    "构建核材料领域垂直专家模型",
    "实现基于RAG的专业知识检索与问答",
    "集成Agent框架支持复杂推理任务",
    "与核电专家合作验证实际应用价值",
    "提升科研工作效率50%以上"
])

# 4. 系统架构
add_two_column_slide(prs, "系统架构",
    "核心技术栈",
    [
        "基座模型: Qwen3.5-2B/9B",
        "训练框架: SWIFT + LoRA",
        "向量化: BGE-M3",
        "知识存储: Milvus + Neo4j",
        "Agent: Plan-Act双模式"
    ],
    "技术特点",
    [
        "领域自适应预训练(DAPT)",
        "知识图谱+向量混合检索",
        "链式思考推理(CoT)",
        "可扩展的工具调用能力",
        "私有化部署保障安全"
    ])

# 5. 项目进展
add_progress_slide(prs)

# 6. 当前成果
add_content_slide(prs, "当前成果", [
    "✓ 完成3000+条核材料领域QA数据清洗",
    "✓ 完成Qwen3.5-2B模型DAPT训练 (1890 steps)",
    "✓ 评估得分: 平均71.2/100 (核材料相关性问题)",
    "✓ 搭建RAG知识检索系统 (Milvus + Neo4j)",
    "✓ 实现Agent Plan-Act推理框架"
])

# 7. Demo演示
add_content_slide(prs, "系统演示", [
    "场景1: 专业问答 - '核材料辐射损伤机制是什么?'",
    "场景2: 知识推理 - 结合文献进行链式分析",
    "场景3: 核电专家协作 - 辅助研究报告撰写",
    "演示地址: http://localhost:8000 (Web服务)"
])

# 8. 面临的困难
add_content_slide(prs, "当前困难与挑战", [
    "高质量领域数据获取成本高",
    "模型在特定任务上仍有知识盲区",
    "RAG检索精度需进一步优化",
    "与核电专家的深度合作还在推进中",
    "算力资源有限，大模型训练受限"
])

# 9. 下一步计划
add_content_slide(prs, "下一步工作计划", [
    "扩大核材料专业数据规模 (目标: 10000+条)",
    "引入更多专家进行模型评估与反馈",
    "优化RAG检索算法，提升准确率",
    "与核电设计院/研究院建立合作",
    "探索多模态能力 (图表理解)",
    "计划6个月内完成实用化部署"
])

# 10. 总结与展望
add_two_column_slide(prs, "项目亮点与价值",
    "技术创新点",
    [
        "DAPT领域自适应预训练",
        "知识图谱增强的RAG系统",
        "Agent驱动的智能推理框架",
        "私有化部署的领域专家助手"
    ],
    "应用价值",
    [
        "提升核材料研究效率",
        "降低专业知识获取门槛",
        "辅助核电工程师日常工作",
        "推动AI4S在核电领域落地",
        "具有广阔的商业前景"
    ])

# 11. 致谢/讨论
add_title_slide(prs, "谢谢!", "欢迎各位老师同学批评指正")

# 保存文件
output_path = "output/核材料专家模型项目汇报.pptx"
os.makedirs("output", exist_ok=True)
prs.save(output_path)
print(f"PPT已生成: {os.path.abspath(output_path)}")
